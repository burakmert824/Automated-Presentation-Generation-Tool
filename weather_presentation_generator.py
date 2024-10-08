
from selenium import webdriver
from selenium.webdriver.common.by import By
from selenium.webdriver.support.ui import WebDriverWait
from selenium.webdriver.support import expected_conditions as EC
import time
from pptx import Presentation
from pptx.util import Inches
import os


image_files = []

# Create a directory if it does not exist
def create_directory(directory):
    if not os.path.exists(directory):
        os.makedirs(directory)

# Find a web element using XPath in the given WebDriver instance.
def find_element_by_xpath(driver, xpath, timeout=10, wait_after=0):
    """
    Find a web element using XPath in the given WebDriver instance.
    Wait for loading the element on the page
    
    Args:
    - driver (webdriver): The WebDriver instance (e.g., ChromeDriver).
    - xpath (str): XPath expression to locate the element.
    - timeout (int): Maximum time to wait for the element to be found (default: 10 seconds).
    
    Returns:
    - WebElement: The found WebElement object if found, else None.
    """
    try:
        # Wait until the element identified by XPath appears
        wait = WebDriverWait(driver, timeout)
        element = wait.until(EC.visibility_of_element_located((By.XPATH, xpath)))
        
        # Wait for 1 second to ensure the element is fully loaded
        if(wait_after):
            time.sleep(wait_after)  

        return element
    
    except Exception as e:
        print(f"Error: {str(e)}")
        return None



# Example usage:
url = 'https://www.accuweather.com/tr/tr/izmit/318719/weather-forecast/318719'

# Initialize WebDriver (e.g., Chrome)
driver = webdriver.Chrome()

# Navigate to the URL
driver.get(url)
time.sleep(1)  # Wait page is loading


####### Closing the cookie banner #######
COOKIE_XPATH = '//*[@id="privacy-policy-banner"]/div/div'
try:
    accept_button = find_element_by_xpath(driver, COOKIE_XPATH)
    accept_button.click()    
except Exception as e:
    print(f"Privacy policy banner not found or already closed. Error: {str(e)}")
#########################################



## get xpatsh off elements and take screenshots of the elements
xpaths_and_elements = {
    'CURRENT_WEATHER' :'/html/body/div/div[7]/div[1]/div[1]/a[1]',
    'NEXT_DAYS_WEATHER' : '/html/body/div/div[7]/div[1]/div[1]/div[2]',
}

# note: initialize element helps for next days weather's png's data download
element = None
SCREENSHOT_DIR = 'screenshots'
create_directory(SCREENSHOT_DIR)

for name , xpath in xpaths_and_elements.items():
    try:
        print(f"Finding element: {name}")
        # Find element by XPath
        element = find_element_by_xpath(driver, xpath,wait_after=0.5)
        image_path = os.path.join(SCREENSHOT_DIR,f'{name}.png')
        # Take screenshot of the current weather element
        element.screenshot(image_path)
        image_files.append(image_path)
        
    except Exception as e:
        print(f"Screenshot failed: {str(e)}")



# Close the browser
driver.quit()


################## Creating Presentation ####################

PRESENTATION_FOLDER = "presentations"
create_directory(PRESENTATION_FOLDER)

# create presentation
prs = Presentation()

# Add title and subtitle to the presentation
lyt = prs.slide_layouts[0] # choose layout style
slide = prs.slides.add_slide(lyt) # add slide 
title = slide.shapes.title # assign a title
subtitle = slide.placeholders[1] # assign a subtitle

title.text = "Hava Durumu" # set title
subtitle.text = "İzmit" # set subtitle


for png_file in image_files:
    lyt_img = prs.slide_layouts[6] # choose layout style
    slide = prs.slides.add_slide(lyt_img) # add slide 
    left = Inches(0.5) # set left margin
    top = Inches(0.5) # set top margin
    pic = slide.shapes.add_picture(png_file, left, top) # add picture to slide



prs.save(os.path.join(PRESENTATION_FOLDER,"weather_presentation.pptx")) # save the presentation
