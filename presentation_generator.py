from pptx import Presentation
from pptx.util import Inches
import os

# Get all PNG files in the folder
def get_png_files(folder_path):
    png_files = []
    try:
        # Iterate over all files in the folder
        for filename in os.listdir(folder_path):
            file_path = os.path.join(folder_path, filename)
            
            # Check if the file is a PNG file
            if filename.lower().endswith('.png') and os.path.isfile(file_path):
                png_files.append(file_path)
                
    except Exception as e:
        print(f"Error while fetching PNG files: {str(e)}")
    
    return png_files

# Create a directory if it does not exist
def create_directory(directory):
    if not os.path.exists(directory):
        os.makedirs(directory)

PRESENTATION_FOLDER = "presentations"
create_directory(PRESENTATION_FOLDER)

# create presentation
prs = Presentation()

# Add title and subtitle to the presentation
lyt = prs.slide_layouts[0] # choose layout style
slide = prs.slides.add_slide(lyt) # add slide 
title = slide.shapes.title # assign a title
subtitle = slide.placeholders[1] # assign a subtitle

title.text = "Hava Durumu" # set title
subtitle.text = "İzmit" # set subtitle

# read png files from the directory and add them to the presentation
SCREENSHOTS_PATH = "screenshots"
png_files = get_png_files(SCREENSHOTS_PATH)
for png_file in png_files:
    lyt_img = prs.slide_layouts[6] # choose layout style
    slide = prs.slides.add_slide(lyt_img) # add slide 
    left = Inches(0.5) # set left margin
    top = Inches(0.5) # set top margin
    pic = slide.shapes.add_picture(png_file, left, top) # add picture to slide



prs.save(os.path.join(PRESENTATION_FOLDER,"weather_presentation.pptx")) # save the presentation
